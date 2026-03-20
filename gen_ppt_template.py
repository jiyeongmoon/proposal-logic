from pptx import Presentation
from pptx.util import Inches, Pt
import os

# Define template path
template_dir = r"c:\Users\user\공간환경계획연구실 Dropbox\04_Knowledge_Base\00_Obsidian\moon\02_데이터_및_스크립트\02_자동화_스크립트\proposal_system\templates"
template_path = os.path.join(template_dir, "standard_temp.pptx")

# Ensure directory exists
os.makedirs(template_dir, exist_ok=True)

def create_sample_template():
    prs = Presentation()
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "{{PROJECT_NAME}}"
    subtitle.text = "사업 제안 발표 자료\n작성일: {{DATE}}"
    
    # Slide 2: Vision Slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    p = tf.add_paragraph()
    p.text = "1. 사업의 비전"
    p.font.bold = True
    p.font.size = Pt(32)
    
    p2 = tf.add_paragraph()
    p2.text = "{{VISION}}"
    p2.font.size = Pt(24)
    
    # Save the presentation
    prs.save(template_path)
    print(f"Sample PPT template created at: {template_path}")

if __name__ == "__main__":
    create_sample_template()
