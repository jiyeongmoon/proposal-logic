from pptx import Presentation
import os

def replace_text_in_ppt(template_path, output_path, replacements):
    """
    PPT 템플릿 내의 {{TAG}}를 실제 텍스트로 교체함.
    replacements: {'TAG_NAME': 'Replacement Text', ...}
    """
    if not os.path.exists(template_path):
        return False, f"템플릿 파일을 찾을 수 없습니다: {template_path}"

    try:
        prs = Presentation(template_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for tag, text in replacements.items():
                            placeholder = f"{{{{{tag}}}}}"  # {{TAG}} 형식
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, str(text))
                                
        prs.save(output_path)
        return True, output_path
    except Exception as e:
        return False, str(e)

if __name__ == "__main__":
    # Test (이 스크립트는 실제 PPT 템플릿이 있어야 작동함)
    test_replacements = {"PROJECT_NAME": "테스트 프로젝트", "VISION": "혁신적인 도시재생"}
    # success, msg = replace_text_in_ppt("template.pptx", "output.pptx", test_replacements)
    # print(f"Success: {success}, Message: {msg}")
