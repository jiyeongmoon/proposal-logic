import streamlit as st
import streamlit.components.v1 as components
import json
import os
import sys
import glob
import pandas as pd
import re
from io import BytesIO
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
import tkinter as tk
from tkinter import filedialog

# 프로젝트 루트 경로 추가
current_dir = os.path.dirname(os.path.abspath(__file__))
if current_dir not in sys.path:
    sys.path.append(current_dir)

from config import APP_TITLE, APP_SUBTITLE, FOLDER_STRUCTURES, BIZ_CODE_DISPLAY, TEMPLATE_DIR, PROPOSAL_BASE_DIR
from utils.folder_logic import create_project_folders
from utils.quant_logic import match_personnel, load_master_db
from utils.ai_assistant import generate_proposal_draft, get_suggested_content_list
from utils.ppt_engine import replace_text_in_ppt
from utils.doc_parser import parse_docx_to_markdown
from utils.gemini_engine import analyze_rfp_with_gemini, get_available_models, refine_logic_with_feedback
from utils.semantic_validator import validate_logic_block
from utils.ui_config import load_css
import utils.prompts as prompts

# --- Page Config ---
st.set_page_config(page_title=APP_TITLE, page_icon="🏢", layout="wide")

def init_session_state():
    """앱에서 사용하는 모든 세션 상태 변수들을 한 곳에서 초기화함"""
    defaults = {
        'current_project_path': None,
        'current_project_name': None,
        'job_type': list(FOLDER_STRUCTURES.keys())[0],
        'biz_code': "UR",
        'proposal_base_dir': PROPOSAL_BASE_DIR,
        'available_models': ["gemini-2.0-flash-exp", "gemini-1.5-flash"],
        'selected_model': "gemini-2.0-flash-exp",
        'nav_index': 0,
        'ai_analysis_res': "", # Background Cache (Guidelines)
        'insight_summary': '', # Background Cache (Insights)
        'tows_result': '', 'tows_data': None,
        'taps_result': '', 'taps_data': None,
        'planner_intent': '',
        'context_files_text': '',
        'show_balloons': False,
        'last_synced_path': None,
        # Widget States (Keys)
        'ai_ta_area': "",
        'ai_ta_area_load': "",
        'insight_ta_area': "",
        'intent_ta_area': "",
        'vision_detail_res': "",
        'final_draft_res': ""
    }
    for key, val in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = val

def find_folder(base_path, pattern):
    """와일드카드를 사용하여 폴더 탐색 (예: 99*)"""
    if not base_path or not os.path.exists(base_path): return None
    matches = glob.glob(os.path.join(base_path, pattern))
    dirs = sorted([m for m in matches if os.path.isdir(m)])
    return dirs[0] if dirs else None

def sync_project_state():
    """프로젝트 선택 시 최신 분석 결과 자동 감지 및 '백그라운드 세션'에만 로드 (UI 위젯 미터치)"""
    if not st.session_state.current_project_path: return
    
    # 1. 가이드라인 로드 (99* 가이드라인 폴더)
    guide_path = find_folder(st.session_state.current_project_path, "99*")
    if guide_path:
        all_mds = sorted([f for f in os.listdir(guide_path) if f.endswith(".md")], reverse=True)
        if all_mds:
            # '[전체분석]' 또는 '[가이드라인]' 프리픽스 파일 검색
            target_files = [f for f in all_mds if f.startswith("[전체분석]") or f.startswith("[가이드라인]")]
            target_file = target_files[0] if target_files else all_mds[0]
            try:
                with open(os.path.join(guide_path, target_file), "r", encoding="utf-8") as f:
                    st.session_state.ai_analysis_res = f.read()
            except: pass

    # 2. 현황요약 (01* 현황분석 폴더)
    survey_path = find_folder(st.session_state.current_project_path, "01*")
    if survey_path:
        files = sorted([f for f in os.listdir(survey_path) if "현황요약" in f and f.endswith(".md")], reverse=True)
        if files:
            try:
                with open(os.path.join(survey_path, files[0]), "r", encoding="utf-8") as f:
                    st.session_state.insight_summary = f.read()
            except: pass
    
    st.session_state.last_synced_path = st.session_state.current_project_path

# --- Utility Functions ---

def select_folder_native():
    """윈도우 표준 폴더 선택창 오픈 (로컬 실행 환경)"""
    root = tk.Tk()
    root.withdraw()
    root.attributes('-topmost', True)
    folder_path = filedialog.askdirectory(initialdir=st.session_state.proposal_base_dir)
    root.destroy()
    return folder_path

def find_folder(base_path, pattern):
    """지정한 패턴(예: '07*')에 매칭되는 하위 폴더 경로 반환"""
    if not base_path or not os.path.exists(base_path): return None
    matches = glob.glob(os.path.join(base_path, pattern))
    dirs = [m for m in matches if os.path.isdir(m)]
    return dirs[0] if dirs else None

def extract_swot_element(content, key_char):
    """현황요약 마크다운에서 S, W, O, T 개별 요소 추출"""
    keyword_map = {'S': '강점', 'W': '약점', 'O': '기회', 'T': '위협'}
    k = keyword_map.get(key_char, "")
    pattern = rf"\|\s*{k}\({key_char}\)\s*\|\s*(.*?)\s*\|"
    match = re.search(pattern, content)
    if match: return match.group(1).strip()
    bullet_pattern = rf"-\s*{k}\({key_char}\):\s*(.*?)(?=\n-|\n\n|$)"
    match = re.search(bullet_pattern, content, re.DOTALL)
    if match: return match.group(1).strip()
    return ""

def extract_ai_section(content, section_name):
    """마크다운 본문에서 특정 섹션 추출"""
    pattern = rf"(?:##\s*\s*|)?\[{section_name}\].*?(?=\n(?:##\s*|)?\[|$)"
    match = re.search(pattern, content, re.DOTALL)
    if match: return match.group(0).strip()
    return f"⚠️ {section_name} 정보를 찾을 수 없습니다."

def extract_gap_analysis(content):
    """현황요약 마크다운에서 🔍 데이터 GAP 분석 섹션 추출"""
    if not content: return ""
    pattern = r"## 🔍 데이터 GAP 분석.*?(?=\n##|$)"
    match = re.search(pattern, content, re.DOTALL)
    if match: return match.group(0).strip()
    # 옛날 버전 또는 이모지 없는 버전 대응
    pattern_alt = r"## 데이터 GAP 분석.*?(?=\n##|$)"
    match_alt = re.search(pattern_alt, content, re.DOTALL)
    if match_alt: return match_alt.group(0).strip()
    return ""

def _extract_text(fpath: str, max_chars: int = 4000) -> str:
    if not os.path.exists(fpath): return ""
    ext = os.path.splitext(fpath)[1].lower()
    try:
        if ext in (".md", ".txt"):
            with open(fpath, "r", encoding="utf-8") as f: return f.read()[:max_chars]
        elif ext == ".docx":
            from docx import Document
            doc = Document(fpath)
            return "\n".join(p.text for p in doc.paragraphs)[:max_chars]
        elif ext == ".pdf":
            import pdfplumber
            with pdfplumber.open(fpath) as pdf:
                pages = [p.extract_text() or "" for p in pdf.pages[:5]]
            return "\n".join(pages)[:max_chars]
        elif ext in (".csv", ".xlsx", ".xls"):
            df = pd.read_csv(fpath) if ext == ".csv" else pd.read_excel(fpath)
            return df.to_string(max_rows=30)[:max_chars]
    except Exception as e: return f"[Error: {e}]"
    return ""

def _assemble_survey_data(base_path):
    SUPPORTED_EXT = {".md", ".txt", ".docx", ".pdf", ".csv", ".xlsx", ".xls"}
    assembled = ""; scan_report = {}
    if not os.path.exists(base_path): return "", {}
    sub_folders = sorted([d for d in os.listdir(base_path) if os.path.isdir(os.path.join(base_path, d)) and not d.startswith('.')])
    for sub in sub_folders:
        sub_path = os.path.join(base_path, sub)
        files = [f for f in os.listdir(sub_path) if os.path.splitext(f)[1].lower() in SUPPORTED_EXT]
        scan_report[sub] = files
        texts = [f"[{f}]\n{_extract_text(os.path.join(sub_path, f))}" for f in files]
        if texts: assembled += f"\n\n=== [{sub}] ===\n" + "\n---\n".join(texts)
    return assembled, scan_report

def extract_json_block(text):
    """마크다운에서 JSON 블록 추출"""
    match = re.search(r"```json\s*(.*?)\s*```", text, re.DOTALL)
    if match:
        try: return json.loads(match.group(1))
        except: return None
    return None

# --- PPT Export Logic ---

def generate_tows_pptx(data):
    prs = Presentation()
    prs.slide_width = Inches(11.69); prs.slide_height = Inches(8.27)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10.6), Inches(0.5)).text_frame
    title_box.text = f"TOWS 전략 분석: {data.get('project_title', '전략 분석')}"
    title_box.paragraphs[0].font.size = Pt(24); title_box.paragraphs[0].font.bold = True
    table = slide.shapes.add_table(3, 3, Inches(0.5), Inches(1), Inches(10.7), Inches(6.5)).table
    sw = data.get('swot', {})
    table.cell(0, 0).text = "TOWS MATRIX"; table.cell(0, 1).text = "STRENGTHS (강점)\n" + "\n".join([f"- {x}" for x in sw.get('strengths', [])])
    table.cell(0, 2).text = "WEAKNESSES (약점)\n" + "\n".join([f"- {x}" for x in sw.get('weaknesses', [])])
    table.cell(1, 0).text = "OPPORTUNITIES (기회)\n" + "\n".join([f"- {x}" for x in sw.get('opportunities', [])])
    table.cell(2, 0).text = "THREATS (위협)\n" + "\n".join([f"- {x}" for x in sw.get('threats', [])])
    def format_strat(key):
        d = data.get(key, {}); txt = f"[{d.get('title', '').upper()}]\n"
        txt += "\n".join([f"• {x}" for x in d.get('bullets', [])]); return txt
    table.cell(1, 1).text = format_strat('so'); table.cell(1, 2).text = format_strat('wo')
    table.cell(2, 1).text = format_strat('st'); table.cell(2, 2).text = format_strat('wt')
    for row in table.rows:
        for cell in row.cells:
            for para in cell.text_frame.paragraphs: para.font.size = Pt(8)
    binary_output = BytesIO(); prs.save(binary_output); return binary_output.getvalue()

def generate_taps_pptx(data):
    prs = Presentation()
    prs.slide_width = Inches(11.69); prs.slide_height = Inches(8.27)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10.6), Inches(0.5)).text_frame
    title_box.text = f"TAPS 전략 체계도: {data.get('project_title', '전략 도출')}"
    title_box.paragraphs[0].font.size = Pt(24); title_box.paragraphs[0].font.bold = True
    axes = data.get('axes', [])
    rows = len(axes) + 1
    table = slide.shapes.add_table(rows, 4, Inches(0.5), Inches(1), Inches(10.7), Inches(6.5)).table
    headers = ["전략축", "현재 문제 (P)", "기본 방향 (D)", "제안 사업 (P)"]
    for i, h in enumerate(headers): table.cell(0, i).text = h
    for i, ax in enumerate(axes):
        table.cell(i+1, 0).text = ax.get('axis_title', '')
        table.cell(i+1, 1).text = ax.get('problem', {}).get('title', '') + "\n" + "\n".join(ax.get('problem', {}).get('bullets', []))
        table.cell(i+1, 2).text = ax.get('direction', {}).get('title', '') + "\n" + "\n".join(ax.get('direction', {}).get('bullets', []))
        table.cell(i+1, 3).text = ax.get('project', {}).get('title', '') + "\n" + "\n".join(ax.get('project', {}).get('bullets', []))
    for row in table.rows:
        for cell in row.cells:
            for para in cell.text_frame.paragraphs: para.font.size = Pt(8)
    binary_output = BytesIO(); prs.save(binary_output); return binary_output.getvalue()

# --- Visualization Helpers ---

def render_tows_html(data):
    if not data: return
    sw = data.get('swot', {})
    C = {'tows_bg': '#F1F5F9', 'str_bg': '#F0FDF4', 'wk_bg': '#FEF2F2', 'opp_bg': '#F5F3FF', 'thr_bg': '#FFF7ED', 'so_bg': '#DCFCE7', 'wo_bg': '#E0E7FF', 'st_bg': '#FFEDD5', 'wt_bg': '#F1F5F9'}
    def items(l): return "".join([f"<li>{x}</li>" for x in l]) if l else "<li>항목 없음</li>"
    def strat_render(key, color):
        d = data.get(key, {}); title = d.get('title', ''); bullets = d.get('bullets', [])
        bullet_html = "".join([f"<div style='margin-bottom:3px;'>• {x}</div>" for x in bullets])
        return f"""<div style="background: {color}; padding: 10px; border-radius: 6px; border: 1px solid #CBD5E1;"><div style="font-weight: 800; color: #1E293B; margin-bottom: 5px; font-size: 0.9rem;">{title}</div><div style="font-size: 0.8rem; line-height: 1.4; color: #334155;">{bullet_html}</div></div>"""
    html = f"""<div style="font-family: 'Pretendard', sans-serif; display: grid; grid-template-columns: 1.2fr 2fr 2fr; gap: 6px; background: #E2E8F0; padding: 6px; border-radius: 8px;"><div style="background: {C['tows_bg']}; padding: 10px; border-radius: 6px; display: flex; align-items: center; justify-content: center; font-weight: 800; color: #475569; font-size: 1rem; border: 1px solid #CBD5E1; text-align:center;">TOWS<br>MATRIX</div><div style="background: {C['str_bg']}; padding: 10px; border-radius: 6px; border-top: 3px solid #22C55E;"><div style="font-size: 0.7rem; color: #16A34A; font-weight: 800; margin-bottom: 3px;">STRENGTHS (강점)</div><ul style="font-size: 0.78rem; padding-left: 15px; margin: 0; color: #166534;">{items(sw.get('strengths'))}</ul></div><div style="background: {C['wk_bg']}; padding: 10px; border-radius: 6px; border-top: 3px solid #EF4444;"><div style="font-size: 0.7rem; color: #DC2626; font-weight: 800; margin-bottom: 3px;">WEAKNESSES (약점)</div><ul style="font-size: 0.78rem; padding-left: 15px; margin: 0; color: #991B1B;">{items(sw.get('weaknesses'))}</ul></div><div style="background: {C['opp_bg']}; padding: 10px; border-radius: 6px; border-left: 3px solid #8B5CF6;"><div style="font-size: 0.7rem; color: #7C3AED; font-weight: 800; margin-bottom: 3px;">OPPORTUNITIES (기회)</div><ul style="font-size: 0.78rem; padding-left: 15px; margin: 0; color: #5B21B6;">{items(sw.get('opportunities'))}</ul></div>{strat_render('so', C['so_bg'])}{strat_render('wo', C['wo_bg'])}<div style="background: {C['thr_bg']}; padding: 10px; border-radius: 6px; border-left: 3px solid #F97316;"><div style="font-size: 0.7rem; color: #EA580C; font-weight: 800; margin-bottom: 3px;">THREATS (위협)</div><ul style="font-size: 0.78rem; padding-left: 15px; margin: 0; color: #9A3412;">{items(sw.get('threats'))}</ul></div>{strat_render('st', C['st_bg'])}{strat_render('wt', C['wt_bg'])}</div>"""
    components.html(html, height=500)

def render_taps_html(data):
    if not data: return
    axes = data.get('axes', [])
    html = f"""<div style="font-family: 'Pretendard', sans-serif; display: flex; flex-direction: column; gap: 10px;">"""
    for ax in axes:
        prob, dire, proj = ax.get('problem',{}), ax.get('direction',{}), ax.get('project',{})
        html += f"""<div style="display: grid; grid-template-columns: 1fr 1fr 1fr; gap: 10px; background: #F8FAFC; padding: 10px; border-radius: 8px; border: 1px solid #E2E8F0;"><div style="grid-column: 1 / 4; font-weight: 800; color: #1E293B; border-bottom: 2px solid #E2E8F0; padding-bottom: 5px; margin-bottom: 5px; font-size: 0.95rem;">{ax.get('axis_title','')}</div><div style="background: #FFF1F2; padding: 8px; border-radius: 6px;"><div style="color: #E11D48; font-weight: 800; font-size: 0.65rem; margin-bottom: 3px;">PROBLEM</div><div style="font-weight: 800; margin-bottom: 5px; font-size: 0.85rem;">{prob.get('title','')}</div><ul style="font-size: 0.78rem; padding-left: 15px; margin: 0; color: #9F1239;">{"".join([f"<li>{x}</li>" for x in prob.get('bullets',[])])}</ul></div><div style="background: #EFF6FF; padding: 8px; border-radius: 6px;"><div style="color: #2563EB; font-weight: 800; font-size: 0.65rem; margin-bottom: 3px;">DIRECTION</div><div style="font-weight: 800; margin-bottom: 5px; font-size: 0.85rem;">{dire.get('title','')}</div><ul style="font-size: 0.78rem; padding-left: 15px; margin: 0; color: #1E40AF;">{"".join([f"<li>{x}</li>" for x in dire.get('bullets',[])])}</ul></div><div style="background: #F0FDF4; padding: 8px; border-radius: 6px; border: 2px solid #22C55E;"><div style="color: #16A34A; font-weight: 800; font-size: 0.65rem; margin-bottom: 3px;">PROJECT (Rank {proj.get('priority','')})</div><div style="font-weight: 800; margin-bottom: 5px; font-size: 0.85rem;">{proj.get('title','')}</div><ul style="font-size: 0.78rem; padding-left: 15px; margin: 0; color: #166534;">{"".join([f"<li>{x}</li>" for x in proj.get('bullets',[])])}</ul></div></div>"""
    html += "</div>"
    components.html(html, height=520, scrolling=True)

# --- Rendering Components ---

@st.dialog("✨ 신규 프로젝트 생성")
def create_project_dialog():
    st.markdown("#### ✨ 1. 신규 프로젝트 정보 입력")
    project_year = datetime.now().strftime("%y")
    c_biz, c_serial = st.columns([0.7, 0.3])
    biz_display = c_biz.selectbox("분류 선택", list(BIZ_CODE_DISPLAY.values()))
    biz_code = biz_display[:2]; serial_num = c_serial.text_input("순번", value="01", max_chars=2)
    project_id = f"{project_year}-{biz_code}{serial_num}"
    name_input = st.text_input("프로젝트명(수행지명) 입력")
    if st.button("🚀 프로젝트 킥오프"):
        if name_input:
            success, result = create_project_folders(project_id, name_input, biz_code, st.session_state.proposal_base_dir)
            if success:
                st.session_state.update({'current_project_path': result, 'current_project_name': name_input, 'job_type': biz_code, 'biz_code': biz_code, 'show_balloons': True})
                sync_project_state(); st.rerun()
            else: st.error(result)

def render_sidebar():
    with st.sidebar:
        st.subheader("🧭 사이드바")
        # --- 프로젝트 탐색 및 생성 (Drill-down) ---
        with st.expander("📂 프로젝트/경로 관리", expanded=False):
            # 1. 경로지정
            st.caption("📍 경로지정")
            col1, col2 = st.columns([0.8, 0.2])
            with col1: st.markdown(f"<code style='font-size:0.75rem;'>.../{os.path.basename(st.session_state.proposal_base_dir)}</code>", unsafe_allow_html=True)
            with col2: 
                if st.button("📁", help="윈도우 창에서 폴더 선택"): 
                    picked = select_folder_native()
                    if picked: st.session_state.proposal_base_dir = picked; st.rerun()
            
            # 2. 신규프로젝트 생성
            st.divider()
            if st.button("➕ 신규 프로젝트 생성", use_container_width=True): create_project_dialog()
            
            # 3. 기존 프로젝트
            st.divider()
            st.caption("📂 기존 프로젝트")
            if os.path.exists(st.session_state.proposal_base_dir):
                items = sorted([d for d in os.listdir(st.session_state.proposal_base_dir) if os.path.isdir(os.path.join(st.session_state.proposal_base_dir, d))], reverse=True)
                sel = st.selectbox("기존 프로젝트 선택 리스트", ["선택하세요"] + items, key="folder_explorer", label_visibility="collapsed")
                if sel != "선택하세요":
                    target = os.path.join(st.session_state.proposal_base_dir, sel)
                    if st.button("📌 선택", use_container_width=True, type="primary", help="이 프로젝트 활성화"):
                        st.session_state.current_project_path = target
                        st.session_state.current_project_name = sel.split("_")[-1]
                        try: st.session_state.biz_code = sel.split("-")[1][:2]
                        except: st.session_state.biz_code = "UR"
                        sync_project_state()
                        if st.session_state.insight_summary and not st.session_state.insight_ta_area: st.session_state.insight_ta_area = st.session_state.insight_summary
                        st.rerun()

        if st.session_state.current_project_path:
            st.markdown(f"<div style='background:#fcf6bd; border-radius:8px; padding:0.5rem 0.8rem; font-size:0.85rem; color:#5465ff;'>📌 <b>{st.session_state.current_project_name}</b></div>", unsafe_allow_html=True)
            if st.button("📂 폴더 열기", use_container_width=True):
                import subprocess
                subprocess.Popen(['explorer', st.session_state.current_project_path])
        st.divider()
        nav_options = ["Phase 1: 가이드라인 스캔", "Phase 2: 사업전략 도출", "🛠️ 툴킷"]
        for i, option in enumerate(nav_options):
            if st.button(option, use_container_width=True, type="primary" if st.session_state.nav_index == i else "secondary"):
                st.session_state.nav_index = i; st.rerun()
        st.divider()
        gemini_api_key = st.text_input("API Key", type="password")
        if st.button("모델 로드") and gemini_api_key:
            st.session_state.available_models = get_available_models(gemini_api_key)
        st.session_state.selected_model = st.selectbox("AI 모델", st.session_state.available_models)
        return gemini_api_key, nav_options[st.session_state.nav_index]

def render_phase1_scanner(api_key):
    st.markdown("### 🔍 가이드라인 자동 스캔 및 분석")
    if not st.session_state.current_project_path: st.warning("프로젝트를 선택하세요."); return
    ref_path = find_folder(st.session_state.current_project_path, "07*") or os.path.join(st.session_state.current_project_path, "07_참고 및 기초자료")
    guide_path = find_folder(st.session_state.current_project_path, "99*") or os.path.join(st.session_state.current_project_path, "99. 가이드라인 폴더")
    os.makedirs(ref_path, exist_ok=True); os.makedirs(guide_path, exist_ok=True)
    
    t_load, t_new = st.tabs(["📂 기존 가이드라인 로드", "📄 새 가이드라인 스캔"])
    
    with t_load:
        all_mds = sorted([f for f in os.listdir(guide_path) if f.endswith(".md")], reverse=True)
        if all_mds:
            sel_md = st.selectbox("파일 선택", all_mds, key="guide_file_select")
            if st.button("📂 선택 가이드라인 분석 로드", type="primary", use_container_width=True, key="load_guide_btn"):
                try:
                    with open(os.path.join(guide_path, sel_md), "r", encoding="utf-8") as f:
                        content = f.read()
                        st.session_state.ai_ta_area_load = content
                        st.session_state.ai_analysis_res = content # 동기화
                    st.success(f"로드 성공: {sel_md}"); st.rerun()
                except Exception as e:
                    st.error(f"파일을 읽는 중 오류 발생: {e}")
        else: st.info("대상 폴더에 Markdown 파일이 없습니다.")
        
        st.text_area("기존 가이드라인 분석 내역 (읽기 전용)", height=500, key="ai_ta_area_load", disabled=True)

    with t_new:
        docx_files = glob.glob(os.path.join(ref_path, "*.docx"))
        if docx_files:
            selected = st.multiselect("변환할 DOCX", [os.path.basename(f) for f in docx_files], default=[os.path.basename(f) for f in docx_files])
            if st.button("Markdown 변환 시작", key="convert_btn"):
                p = st.progress(0); total = len(selected)
                for i, name in enumerate(selected):
                    md_text, _ = parse_docx_to_markdown(os.path.join(ref_path, name))
                    with open(os.path.join(ref_path, f"[MD]_{name.replace('.docx','.md')}"), "w", encoding="utf-8") as m: m.write(md_text)
                    p.progress((i + 1) / total)
                st.success("변환 완료")
        
        st.divider()
        if st.button("🤖 AI 핵심 전략 분석 실행", type="primary", use_container_width=True, key="ai_analyze_btn"):
            if not api_key: st.error("API Key 필수"); return
            all_mds = glob.glob(os.path.join(ref_path, "*.md"))
            if all_mds:
                context = "".join([open(f, "r", encoding="utf-8").read() for f in all_mds])
                with st.status("RFP 심층 분석 중...", expanded=True):
                    scan_prompt = prompts.get_rfp_scan_prompt(context)
                    res, _ = analyze_rfp_with_gemini(api_key, scan_prompt, st.session_state.selected_model)
                    if res:
                        st.session_state.ai_analysis_res = res; st.session_state.ai_ta_area = res
                        out_name = f"[전체분석]_{st.session_state.current_project_name}_{datetime.now().strftime('%m%d_%H%M')}.md"
                        with open(os.path.join(guide_path, out_name), "w", encoding="utf-8") as f: f.write(res)
                        st.success("분석 완료!"); st.rerun()
            else: st.warning("분석할 Markdown 파일이 07. 폴더에 없습니다.")
        
        st.text_area("분석 결과 미리보기 (읽기 전용)", height=400, key="ai_ta_area", disabled=True)

def render_phase2_1_strategy(api_key):
    st.markdown("### 🏹 사업 전략 및 논리고도화")
    if not st.session_state.current_project_path: st.warning("프로젝트 선택 필수"); return
    survey_base = find_folder(st.session_state.current_project_path, "01*") or os.path.join(st.session_state.current_project_path, "01_기초조사 및 현황분석")
    output_base = find_folder(st.session_state.current_project_path, "02*") or os.path.join(st.session_state.current_project_path, "02_기본구상 및 전략")
    os.makedirs(survey_base, exist_ok=True); os.makedirs(output_base, exist_ok=True)
    
    # STEP 1: Insight Summary
    st.markdown("#### :arrow_forward: STEP 1. 데이터 기반 현황 요약")
    t_load, t_scan = st.tabs(["📄 기존 요약문 불러오기", "🤖 신규 AI 생성"])
    with t_load:
        existing = sorted([f for f in os.listdir(survey_base) if "현황요약" in f and f.endswith(".md")], reverse=True)
        if existing:
            sel_md = st.selectbox("파일 선택", existing, key="insight_file_select")
            if st.button("선택 요약 로드", key="insight_load_btn"):
                try:
                    with open(os.path.join(survey_base, sel_md), "r", encoding="utf-8") as f:
                        content = f.read()
                        st.session_state.insight_ta_area = content
                    st.success(f"로드 성공: {sel_md}"); st.rerun()
                except Exception as e:
                    st.error(f"파일을 읽는 중 오류 발생: {e}")
        else: st.info("대상 폴더에 요약 파일이 없습니다.")
    with t_scan:
        if st.button("🤖 현황 요약 생성", type="primary", key="insight_gen_btn"):
            if not api_key: st.error("API Key 필수"); return
            with st.status("데이터 분석 중...", expanded=True):
                assembled, _ = _assemble_survey_data(survey_base)
                guide_for_ai = st.session_state.get('ai_ta_area')
                if not guide_for_ai: guide_for_ai = st.session_state.get('ai_ta_area_load')
                if not guide_for_ai: guide_for_ai = st.session_state.ai_analysis_res
                
                rfp_context = extract_ai_section(guide_for_ai, "사업개요")
                res, _ = analyze_rfp_with_gemini(api_key, prompts.get_insight_prompt(rfp_context, assembled), st.session_state.selected_model)
                if res: 
                    st.session_state.insight_summary = res; st.session_state.insight_ta_area = res
                    with open(os.path.join(survey_base, f"{datetime.now().strftime('%Y%m%d_%H%M')}_현황요약_AI분석.md"), "w", encoding="utf-8") as f: f.write(res)
                    st.success("생성 완료"); st.rerun()
    
    st.text_area("현황요약 결과 편집", height=300, key="insight_ta_area")

    # STEP 1.5: Consultant's Desk
    st.divider(); st.markdown("#### :arrow_forward: STEP 1.5 기획자 의도 주입 및 참고자료")
    
    # GAP 분석 결과 암묵적 노출 (과정 중 참고)
    gap_content = extract_gap_analysis(st.session_state.insight_ta_area)
    if gap_content:
        with st.expander("🔍 데이터 GAP 분석 및 보완 제안 (참고)", expanded=True):
            st.info("현재 분석된 현황과 RFP 가이드라인을 대조한 결과, 아래 자료의 추가 보완을 권장합니다.")
            st.markdown(gap_content)
    
    c_left, c_right = st.columns([2, 1])
    with c_right:
        uploaded_files = st.file_uploader("회의록/참고파일 업로드 (의도 연동)", accept_multiple_files=True)
        if uploaded_files:
            ctx_all = ""
            for f in uploaded_files: ctx_all += f"\n\n[{f.name}]\n" + f.read().decode("utf-8", errors="ignore")
            if ctx_all and ctx_all != st.session_state.context_files_text:
                st.session_state.context_files_text = ctx_all
                combined = (st.session_state.planner_intent + "\n\n[참고 파일 내용]\n" + ctx_all).strip()
                st.session_state.planner_intent = combined; st.session_state.intent_ta_area = combined
                st.info("의도 연동 완료"); st.rerun()
    with c_left:
        st.text_area("기획자 주요 의도 및 참고 맥락", height=180, key="intent_ta_area", placeholder="기획 의도 입력 또는 파일 업로드")

    # STEP 2: Chained Analysis
    st.divider(); st.markdown("#### :arrow_forward: STEP 2. 전략 연쇄 분석 (TOWS → TAPS)")
    if st.button("🚀 TOWS-TAPS 통합 분석 실행", type="primary", use_container_width=True):
        if not api_key: st.error("API Key 필수"); return
        with st.status("전문가 전략 수립 중...", expanded=True) as status:
            guide_final = st.session_state.get('ai_ta_area')
            if not guide_final: guide_final = st.session_state.get('ai_ta_area_load')
            if not guide_final: guide_final = st.session_state.ai_analysis_res
            
            insight_final = st.session_state.get('insight_ta_area') if st.session_state.get('insight_ta_area') else st.session_state.insight_summary
            planner_intent = st.session_state.get('intent_ta_area', '')
            
            s = extract_swot_element(insight_final, 'S'); w = extract_swot_element(insight_final, 'W')
            o = extract_swot_element(insight_final, 'O'); t = extract_swot_element(insight_final, 'T')
            
            rfp_qual = extract_ai_section(guide_final, "정성_평가지표")
            status.update(label="1/2: TOWS 4사분면 교차 전략 도출 중...")
            t_prompt = prompts.get_tows_prompt(rfp_qual, insight_final + "\n" + planner_intent, s, w, o, t, "", planner_intent)
            t_res, _ = analyze_rfp_with_gemini(api_key, t_prompt, st.session_state.selected_model)
            if t_res:
                st.session_state.tows_result = t_res; st.session_state.tows_data = extract_json_block(t_res)
                status.update(label="2/2: TAPS 정성 전략 및 로직 트리 도출 중...")
                a_prompt = prompts.get_taps_prompt(insight_final + "\n" + t_res, "", planner_intent, "", rfp_qual)
                a_res, _ = analyze_rfp_with_gemini(api_key, a_prompt, st.session_state.selected_model)
                if a_res:
                    st.session_state.taps_result = a_res; st.session_state.taps_data = extract_json_block(a_res)
                    ts = datetime.now().strftime("%Y%m%d_%H%M")
                    with open(os.path.join(output_base, f"{ts}_TOWS_전략도출.md"), "w", encoding="utf-8") as f: f.write(t_res)
                    with open(os.path.join(output_base, f"{ts}_TAPS_전략도출.md"), "w", encoding="utf-8") as f: f.write(a_res)
                    st.success("전략 패키지 생성 완료")

    # Matrix Visualization
    c1, c2 = st.columns([1, 1])
    with c1:
        st.subheader("📊 TOWS Matrix")
        if st.session_state.tows_data:
            render_tows_html(st.session_state.tows_data)
            st.download_button("📥 TOWS PPT", generate_tows_pptx(st.session_state.tows_data), "TOWS_Expert.pptx", use_container_width=True)
    with c2:
        st.subheader("🏹 TAPS Strategy Tree")
        if st.session_state.taps_data:
            render_taps_html(st.session_state.taps_data)
            st.download_button("📥 TAPS PPT", generate_taps_pptx(st.session_state.taps_data), "TAPS_Expert.pptx", use_container_width=True)

    # STEP 3: Vision & Detail
    st.divider(); st.markdown("#### :arrow_forward: STEP 3. 비전 및 전략 고도화")
    if st.button("✨ 1줄 비전 / 3단 거버넌스 / 성과지표 도출", type="primary", use_container_width=True):
        if not api_key: st.error("API Key 필수"); return
        with st.status("비전 및 실행 디테일 설계 중...", expanded=True):
            taps_context = st.session_state.get('taps_result', '')
            insight_context = st.session_state.get('insight_ta_area', '')
            intent_context = st.session_state.get('intent_ta_area', '')
            
            if not taps_context: st.warning("Step 2 전략 도출이 선행되어야 합니다."); return
            
            v_prompt = prompts.get_vision_detail_prompt(taps_context, insight_context, intent_context)
            v_res, _ = analyze_rfp_with_gemini(api_key, v_prompt, st.session_state.selected_model)
            if v_res:
                st.session_state.vision_detail_res = v_res
                with open(os.path.join(output_base, f"{datetime.now().strftime('%Y%m%d_%H%M')}_비전및전략고도화.md"), "w", encoding="utf-8") as f: f.write(v_res)
                st.success("비전 및 디테일 도출 완료"); st.rerun()

    if st.session_state.vision_detail_res:
        with st.container(border=True):
            st.markdown("### 📄 비전 및 세부 추진 계획 (Step 3 분석)")
            st.markdown(st.session_state.vision_detail_res)

    # STEP 4: Final Proposal Draft
    st.divider(); st.markdown("#### :arrow_forward: STEP 4. 최종 사업기획 초안 생성")
    if st.button("📄 통합 사업기획 초안 도출 (Final Draft)", type="primary", use_container_width=True):
        if not api_key: st.error("API Key 필수"); return
        with st.status("데이터 집합 및 최종 초안 작성 중...", expanded=True):
            # 데이터 확보 (우선순위: 편집창 내용 > 백그라운드 캐시)
            guide_context = st.session_state.get('ai_ta_area') or st.session_state.get('ai_ta_area_load') or st.session_state.ai_analysis_res
            tows_ctx = st.session_state.tows_result
            taps_ctx = st.session_state.taps_result
            vision_ctx = st.session_state.vision_detail_res
            
            if not (tows_ctx and taps_ctx and vision_ctx):
                st.warning("이전 단계(Step 2, 3)의 전략 도출이 완료되어야 합리적인 초안 생성이 가능합니다."); return

            f_prompt = prompts.get_final_proposal_prompt(
                st.session_state.current_project_name, 
                guide_context, tows_ctx, taps_ctx, vision_ctx
            )
            f_res, _ = analyze_rfp_with_gemini(api_key, f_prompt, st.session_state.selected_model)
            if f_res:
                st.session_state.final_draft_res = f_res
                ts = datetime.now().strftime("%Y%m%d_%H%M")
                out_file = f"{ts}_사업기획초안_행정검토용.md"
                with open(os.path.join(output_base, out_file), "w", encoding="utf-8") as f: f.write(f_res)
                st.success("최종 초안 생성 및 저장 완료"); st.rerun()

    if st.session_state.final_draft_res:
        with st.container(border=True):
            st.markdown("### 📄 최종 사업기획 초안 미리보기")
            st.markdown(st.session_state.final_draft_res)
        st.download_button("📥 최종 초안 다운로드 (.md)", st.session_state.final_draft_res, f"{st.session_state.current_project_name}_최종초안.md", use_container_width=True)

def render_toolkit(api_key):
    st.markdown("### 🛠️ 비즈니스 프레임워크 툴킷")
    source = st.session_state.get('insight_ta_area') if st.session_state.get('insight_ta_area') else st.session_state.insight_summary
    ctx = st.text_area("워크 스페이스 데이터", value=source, height=150)
    fw = st.selectbox("프레임워크", ["AS-IS / TO-BE", "PEST", "Stakeholder"])
    if st.button("⚙️ 프레임워크 렌더링"):
        with st.spinner("AI 분석 중..."):
            res, _ = analyze_rfp_with_gemini(api_key, f"{ctx}\n\n 위 데이터를 기반으로 {fw} 분석 수행", st.session_state.selected_model)
            if res: st.session_state.tk_result = res; st.success("완료")
    if st.session_state.get('tk_result'): st.markdown(st.session_state.tk_result)

# --- Main App ---

if __name__ == "__main__":
    init_session_state()
    load_css(); api_key, page = render_sidebar()
    st.title(APP_TITLE)
    if st.session_state.show_balloons: st.balloons(); st.session_state.show_balloons = False
    if page == "Phase 1: 가이드라인 스캔": render_phase1_scanner(api_key)
    elif page == "Phase 2: 사업전략 도출": render_phase2_1_strategy(api_key)
    elif page == "🛠️ 툴킷": render_toolkit(api_key)
    st.divider(); st.caption("© 2026 나래공간환경연구소 - Advanced Proposal System")
